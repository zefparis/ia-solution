import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
import base64
import requests
import io

def create_powerpoint_from_html(html_path, output_path, language="fr"):
    """
    Crée une présentation PowerPoint à partir du contenu HTML
    
    Args:
        html_path (str): Chemin vers le fichier HTML
        output_path (str): Chemin du fichier PowerPoint à générer
        language (str): Code de langue ("fr" ou "en")
    """
    # Lire le contenu HTML
    with open(html_path, 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Récupérer toutes les slides
    slides = soup.select('.slide')
    
    # Créer une nouvelle présentation
    prs = Presentation()
    
    # Personnaliser la taille des diapositives (format 16:9)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Génération des slides
    for slide_html in slides:
        # Récupérer l'identifiant de la slide
        slide_id = slide_html.get('id', '')
        
        # Pour la slide 10, appliquer un arrière-plan spécial (dégradé)
        if slide_id == 'slide10':
            slide_layout = prs.slide_layouts[5]  # Titre et contenu
            slide = prs.slides.add_slide(slide_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(30, 58, 138)  # Couleur similaire au dégradé
        else:
            slide_layout = prs.slide_layouts[5]  # Titre et contenu
            slide = prs.slides.add_slide(slide_layout)
        
        # Titre de la slide
        title_text = slide_html.select_one('h2')
        if title_text:
            title = slide.shapes.title
            title.text = title_text.get_text(strip=True)
            title_paragraph = title.text_frame.paragraphs[0]
            title_paragraph.alignment = PP_ALIGN.CENTER
            
            # Style du titre
            title_run = title_paragraph.runs[0]
            title_run.font.size = Pt(40)
            title_run.font.bold = True
            
            # Couleur blanche pour le titre de la slide 10
            if slide_id == 'slide10':
                title_run.font.color.rgb = RGBColor(255, 255, 255)
        
        # Zone de contenu
        content_box = slide.placeholders[1]
        content_frame = content_box.text_frame
        content_frame.clear()  # Supprimer le texte par défaut
        
        # Extraire et traiter les sections de contenu
        
        # 1. Extraire les listes
        feature_lists = slide_html.select('.feature-list')
        for feature_list in feature_lists:
            list_items = feature_list.select('li')
            for item in list_items:
                # Récupérer le texte en gras (titre)
                strong_text = item.select_one('strong')
                if strong_text:
                    p = content_frame.add_paragraph()
                    p.level = 0
                    
                    # Ajouter le titre en gras
                    run = p.add_run()
                    run.text = "• " + strong_text.get_text(strip=True)
                    run.font.bold = True
                    if slide_id == 'slide10':
                        run.font.color.rgb = RGBColor(255, 255, 255)
                    
                    # Ajouter la description
                    description = item.get_text().replace(strong_text.get_text(), '').strip()
                    if description:
                        p = content_frame.add_paragraph()
                        p.level = 1
                        run = p.add_run()
                        run.text = description
                        if slide_id == 'slide10':
                            run.font.color.rgb = RGBColor(255, 255, 255)
        
        # 2. Traiter les témoignages
        testimonials = slide_html.select('.testimonial')
        for testimonial in testimonials:
            p = content_frame.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            
            # Ajouter un espace avant le témoignage
            space_p = content_frame.add_paragraph()
            space_p.text = " "
            
            # Ajouter le témoignage
            run = p.add_run()
            
            # Récupérer l'auteur
            author = testimonial.select_one('.author')
            author_text = ""
            if author:
                author_text = author.get_text(strip=True)
                # Retirer le texte de l'auteur du témoignage
                testimonial_text = testimonial.get_text().replace(author_text, '').strip()
            else:
                testimonial_text = testimonial.get_text(strip=True)
            
            run.text = f'"{testimonial_text}"'
            run.font.italic = True
            
            # Ajouter l'auteur sur une nouvelle ligne
            if author_text:
                author_p = content_frame.add_paragraph()
                author_p.alignment = PP_ALIGN.CENTER
                author_run = author_p.add_run()
                author_run.text = author_text
                author_run.font.bold = True
        
        # 3. Traiter les tableaux de prix (slide 9)
        if slide_id == 'slide9':
            pricing_columns = slide_html.select('.pricing-column')
            
            # Ajouter un espace avant les tarifs
            space_p = content_frame.add_paragraph()
            space_p.text = " "
            
            for i, pricing in enumerate(pricing_columns):
                header = pricing.select_one('.pricing-header')
                price = pricing.select_one('.pricing-price')
                features = pricing.select('.pricing-features li')
                
                p = content_frame.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                
                # Titre du plan
                run = p.add_run()
                run.text = header.get_text(strip=True) if header else f"Plan {i+1}"
                run.font.bold = True
                run.font.size = Pt(28)
                
                # Prix
                if price:
                    price_p = content_frame.add_paragraph()
                    price_p.alignment = PP_ALIGN.CENTER
                    price_run = price_p.add_run()
                    price_run.text = price.get_text(strip=True).replace('\n', ' - ')
                    price_run.font.size = Pt(24)
                
                # Stockage
                storage = pricing.select_one('.pricing-content > div:nth-of-type(2)')
                if storage:
                    storage_p = content_frame.add_paragraph()
                    storage_p.alignment = PP_ALIGN.CENTER
                    storage_run = storage_p.add_run()
                    storage_run.text = storage.get_text(strip=True)
                
                # Fonctionnalités
                for feature in features:
                    feature_p = content_frame.add_paragraph()
                    feature_p.alignment = PP_ALIGN.CENTER
                    feature_run = feature_p.add_run()
                    feature_run.text = "✓ " + feature.get_text(strip=True)
                
                # Ajouter un espacement entre les plans
                if i < len(pricing_columns) - 1:
                    separator_p = content_frame.add_paragraph()
                    separator_p.text = "---"
                    separator_p.alignment = PP_ALIGN.CENTER
        
        # 4. Traiter le contenu en deux colonnes
        two_columns = slide_html.select('.two-columns')
        for columns in two_columns:
            column_items = columns.select('.column')
            
            # Récupérer les images dans les colonnes
            images = slide_html.select('img')
            
            # Traiter les images si présentes
            if images:
                for img in images:
                    img_src = img.get('src', '')
                    if img_src and img_src.startswith('/presentation/images/'):
                        # Composer le chemin complet de l'image
                        img_path = img_src.replace('/presentation/images/', 'presentation/images/')
                        if os.path.exists(img_path):
                            # Calculer la position appropriée pour l'image
                            left = Inches(8)  # Position à droite de la slide
                            top = Inches(2.5)  # Position vers le milieu
                            width = Inches(6)  # Largeur raisonnable
                            
                            # Ajouter l'image à la slide
                            slide.shapes.add_picture(img_path, left, top, width=width)
        
        # 5. Traiter les éléments highlight-box
        highlight_boxes = slide_html.select('.highlight-box')
        for box in highlight_boxes:
            p = content_frame.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            
            # Ajouter un espace avant la boîte
            space_p = content_frame.add_paragraph()
            space_p.text = " "
            
            run = p.add_run()
            run.text = box.get_text(strip=True)
            run.font.bold = True
            run.font.size = Pt(20)
            
            # Couleur spéciale pour les boîtes highlight
            run.font.color.rgb = RGBColor(251, 140, 0)  # Orange
    
    # Sauvegarder la présentation
    prs.save(output_path)
    
    return output_path


if __name__ == "__main__":
    # Vérifier si le dossier 'downloads' existe, sinon le créer
    if not os.path.exists('downloads'):
        os.makedirs('downloads')
    
    # Créer les versions française et anglaise
    fr_output = create_powerpoint_from_html(
        "presentation/preview.html", 
        "downloads/presentation-ia-solution-fr.pptx",
        "fr"
    )
    
    en_output = create_powerpoint_from_html(
        "presentation/preview_en.html", 
        "downloads/presentation-ia-solution-en.pptx", 
        "en"
    )
    
    print(f"Présentation PowerPoint française créée: {fr_output}")
    print(f"Présentation PowerPoint anglaise créée: {en_output}")